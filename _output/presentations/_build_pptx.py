from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (light theme) ─────────────────────────────────────────────────────
TEXT     = RGBColor(0x1F, 0x23, 0x28)   # near-black
MUTED    = RGBColor(0x63, 0x6E, 0x7B)   # medium gray
SURFACE  = RGBColor(0xF6, 0xF8, 0xFA)   # very light gray
BORDER   = RGBColor(0xD0, 0xD7, 0xDE)   # light border

GREEN    = RGBColor(0x1A, 0x7F, 0x37)
GREEN_BG = RGBColor(0xDC, 0xFE, 0xE7)
GREEN_BD = RGBColor(0x4A, 0xC2, 0x6B)

YELLOW   = RGBColor(0x7D, 0x4E, 0x00)
YELLOW_BG= RGBColor(0xFF, 0xF8, 0xC5)
YELLOW_BD= RGBColor(0xD4, 0xA7, 0x2C)

RED      = RGBColor(0xCF, 0x22, 0x2E)
RED_BG   = RGBColor(0xFF, 0xEB, 0xE9)
RED_BD   = RGBColor(0xFF, 0x81, 0x82)

BLUE     = RGBColor(0x09, 0x69, 0xDA)
BLUE_BG  = RGBColor(0xDF, 0xEF, 0xFF)
BLUE_BD  = RGBColor(0x54, 0xAE, 0xFF)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    """Blank slide — no background fill, inherits template background."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def add_rect(slide, x, y, w, h, fill=None, border=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = border_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(12), bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def add_flow_step(slide, name, note, x, y, w, h=Inches(0.72),
                  fill=SURFACE, border=BORDER, name_color=TEXT):
    add_rect(slide, x, y, w, h, fill=fill, border=border)
    add_text(slide, name, x, y + Inches(0.1), w, Inches(0.3),
             size=Pt(12), bold=True, color=name_color, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x, y + Inches(0.38), w, Inches(0.28),
                 size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════
prs = new_prs()


# ─────────────────────────────────────────────────────────── SLIDE 1 — Processes
slide = blank_slide(prs)

add_text(slide, "Hur processen ser ut — och hur den borde se ut",
         Inches(0.7), Inches(0.4), Inches(11), Inches(0.65),
         size=Pt(34), bold=True, color=TEXT)
add_rect(slide, Inches(0.7), Inches(1.18), Inches(11.9), Inches(0.02),
         fill=BORDER, border=None)

# ── Row labels ────────────────────────────────────────────────────────────────
add_text(slide, "NU", Inches(0.7), Inches(1.42), Inches(0.9), Inches(0.3),
         size=Pt(10), bold=True, color=RED)
add_text(slide, "BÄTTRE", Inches(0.7), Inches(4.2), Inches(1.2), Inches(0.3),
         size=Pt(10), bold=True, color=GREEN)

# Horizontal separator between rows
add_rect(slide, Inches(0.7), Inches(3.82), Inches(11.9), Inches(0.02),
         fill=BORDER, border=None)

# ── Helper: flow box ──────────────────────────────────────────────────────────
def flow_box(slide, label, sub, x, y, w=Inches(2.2), h=Inches(1.4),
             fill=SURFACE, border=BORDER, lc=TEXT):
    add_rect(slide, x, y, w, h, fill=fill, border=border, border_width=Pt(1.5))
    add_text(slide, label, x, y + Inches(0.38), w, Inches(0.42),
             size=Pt(15), bold=True, color=lc, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, x, y + Inches(0.82), w, Inches(0.38),
                 size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

def arrow(slide, x, y, symbol="→", col=MUTED):
    add_text(slide, symbol, x, y, Inches(0.55), Inches(0.4),
             size=Pt(22), bold=True, color=col, align=PP_ALIGN.CENTER)

# ── NOW row ───────────────────────────────────────────────────────────────────
bw = Inches(2.2)
bh = Inches(1.5)
aw = Inches(0.55)
now_y = Inches(1.72)
now_x = Inches(1.75)
spacing = bw + aw

flow_box(slide, "Krav",        "",                    now_x,               now_y, bw, bh)
arrow(slide, now_x + bw,       now_y + Inches(0.55))
flow_box(slide, "Design",      "",                    now_x + spacing,     now_y, bw, bh)
arrow(slide, now_x + spacing*2,now_y + Inches(0.55))
flow_box(slide, "Utveckling",  "",                    now_x + spacing*2,   now_y, bw, bh)
arrow(slide, now_x + spacing*3,now_y + Inches(0.55), col=RED)
flow_box(slide, "Reaktiv UI/UX", "Efterhandsgranskning",
         now_x + spacing*3,    now_y, bw, bh,
         fill=RED_BG, border=RED_BD, lc=RED)

# Problem callout
add_text(slide, "⚠  UI/UX kommer för sent — dyrt att rätta",
         now_x + spacing*3 - Inches(0.3), now_y + bh + Inches(0.12),
         bw + Inches(0.6), Inches(0.3),
         size=Pt(10), color=RED, align=PP_ALIGN.CENTER)

# ── BETTER row ────────────────────────────────────────────────────────────────
bet_y = Inches(4.3)
bet_x = Inches(1.75)

# Wider middle block for the combined UI/UX <-> Design
mid_w = Inches(3.6)
mid_x = bet_x + spacing

flow_box(slide, "Krav",  "", bet_x, bet_y, bw, bh)
arrow(slide, bet_x + bw, bet_y + Inches(0.55))

# Combined UI/UX <-> Design block
add_rect(slide, mid_x, bet_y, mid_w, bh, fill=BLUE_BG, border=BLUE_BD, border_width=Pt(2))
add_text(slide, "UI/UX  ↔  Design", mid_x, bet_y + Inches(0.32), mid_w, Inches(0.45),
         size=Pt(15), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Iterativt — innan kod skrivs", mid_x, bet_y + Inches(0.78), mid_w, Inches(0.32),
         size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

arrow(slide, mid_x + mid_w, bet_y + Inches(0.55), col=GREEN)
dev_x = mid_x + mid_w + aw
flow_box(slide, "Utveckling", "Tydlig spec", dev_x, bet_y, bw, bh,
         fill=GREEN_BG, border=GREEN_BD, lc=GREEN)

# Benefit callout
add_text(slide, "✓  Alla frågor besvarade innan första raden kod",
         mid_x - Inches(0.2), bet_y + bh + Inches(0.12),
         mid_w + bw + Inches(0.8), Inches(0.3),
         size=Pt(10), color=GREEN, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────── SLIDE 2
slide = blank_slide(prs)

add_text(slide, "Ju senare felet hittas,\ndesto dyrare att rätta",
         Inches(0.7), Inches(0.4), Inches(8), Inches(1.2),
         size=Pt(34), bold=True, color=TEXT)
add_text(slide, "En ändring i designfasen kostar bråkdelen av en ändring i kod",
         Inches(0.7), Inches(1.6), Inches(9), Inches(0.35),
         size=Pt(13), color=MUTED)

# Horizontal rule under subtitle
add_rect(slide, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.02),
         fill=BORDER, border=None)

# ── Bar chart ─────────────────────────────────────────────────────────────────
chart_x  = Inches(0.9)
chart_y  = Inches(2.3)
bar_w    = Inches(1.8)
gap      = Inches(1.0)
max_h    = Inches(3.6)
baseline = chart_y + max_h

bars = [
    ("1×",   "I design",     "Flytta en ruta i prototypen",       0.14, GREEN_BG,  GREEN_BD,  GREEN),
    ("10×",  "I kod",        "Refaktorera komponenter",            0.45, YELLOW_BG, YELLOW_BD, YELLOW),
    ("100×", "I produktion", "Felanalys, hotfix, regressionstest", 1.0,  RED_BG,    RED_BD,    RED),
]

for i, (mult, label, sub, ratio, bg, bd, fc) in enumerate(bars):
    bx = chart_x + i * (bar_w + gap)
    bh = Emu(int(max_h * ratio))
    by = baseline - bh
    add_rect(slide, bx, by, bar_w, bh, fill=bg, border=bd, border_width=Pt(1.5))
    add_text(slide, mult,  bx, by + Inches(0.1), bar_w, Inches(0.5),
             size=Pt(24), bold=True, color=fc, align=PP_ALIGN.CENTER)
    add_text(slide, label, bx, baseline + Inches(0.1), bar_w, Inches(0.3),
             size=Pt(12), bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, sub,   bx, baseline + Inches(0.42), bar_w, Inches(0.4),
             size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)

# ── Insight cards ─────────────────────────────────────────────────────────────
cx  = Inches(7.3)
cy  = Inches(2.3)
cw  = Inches(5.6)
cg  = Inches(0.22)
ch  = Inches(1.12)

insights = [
    (RED,   RED_BG,   RED_BD,   "Utan UX-first",
     "Utvecklare tolkar vaga krav. Stakeholders ser resultatet och\nsäger \"det var inte det vi menade\" — arbetet börjar om."),
    (RED,   RED_BG,   RED_BD,   "Omarbetning är det största slöseriet",
     "Kod som kastas kostar mer än kod som aldrig skrivs.\nVarje omarbetningssprint är en sprint utan nytt värde."),
    (GREEN, GREEN_BG, GREEN_BD, "Lösningen: fånga det i design",
     "Prototyper är billiga att ändra. En godkänd design innan\nkod skrivs eliminerar den dyraste kategorin av fel."),
]

for i, (fc, bg, bd, title, body) in enumerate(insights):
    iy = cy + i * (ch + cg)
    add_rect(slide, cx, iy, cw, ch, fill=bg, border=bd, border_width=Pt(1.5))
    add_text(slide, title, cx + Inches(0.18), iy + Inches(0.1),
             cw - Inches(0.3), Inches(0.3),
             size=Pt(12), bold=True, color=fc)
    add_text(slide, body,  cx + Inches(0.18), iy + Inches(0.42),
             cw - Inches(0.3), Inches(0.62),
             size=Pt(10), color=TEXT)


# ─────────────────────────────────────────────────────────── SLIDE 3
slide = blank_slide(prs)

add_text(slide, "Vad en godkänd design ger utvecklaren",
         Inches(0.7), Inches(0.4), Inches(11), Inches(0.65),
         size=Pt(34), bold=True, color=TEXT)
add_text(slide, "Svaret på alla frågor — innan de ens hinner ställas",
         Inches(0.7), Inches(1.08), Inches(9), Inches(0.35),
         size=Pt(13), color=MUTED)
add_rect(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.02),
         fill=BORDER, border=None)

# Left: deliverable cards
lx = Inches(0.7)
ly = Inches(1.72)
lw = Inches(5.6)
lh = Inches(1.52)
lg = Inches(0.18)

deliverables = [
    (BLUE,  "VISUELLT",  "Annoterade skärmar",
     ["Alla skärmar och tillstånd: tom, laddning, fel, ifylld",
      "Exakta mått, typsnitt och färger från designsystemet",
      "Responsivt beteende specificerat"]),
    (BLUE,  "BETEENDE",  "Interaktionsspec",
     ["Vad händer vid varje klick, hover och scroll",
      "Felhantering och kantfall definierade",
      "Övergångar och laddningstillstånd"]),
    (BLUE,  "TEKNIK",    "Datakrav & komponenter",
     ["Vilka API-anrop skärmen behöver",
      "Existerande komponenter som återanvänds",
      "Nya komponenter med spec bifogad"]),
]

for i, (fc, lbl, title, items) in enumerate(deliverables):
    y = ly + i * (lh + lg)
    add_rect(slide, lx, y, lw, lh, fill=SURFACE, border=BORDER, border_width=Pt(1))
    add_rect(slide, lx, y, Inches(0.06), lh, fill=fc, border=None)  # left accent stripe
    add_text(slide, lbl,   lx + Inches(0.2), y + Inches(0.1),  lw - Inches(0.3), Inches(0.25),
             size=Pt(9), bold=True, color=fc)
    add_text(slide, title, lx + Inches(0.2), y + Inches(0.35), lw - Inches(0.3), Inches(0.3),
             size=Pt(13), bold=True, color=TEXT)
    body = "\n".join(f"–  {l}" for l in items)
    add_text(slide, body,  lx + Inches(0.2), y + Inches(0.68), lw - Inches(0.3), lh - Inches(0.75),
             size=Pt(10), color=MUTED)

# Right: outcome list + statement
rx  = Inches(7.0)
ry  = Inches(1.72)
rw  = Inches(5.9)
rh  = Inches(0.88)
rg  = Inches(0.16)

outcomes = [
    (GREEN, "Utvecklaren kan köra direkt",
     "Ingen tid spenderas på att tolka krav eller vänta på svar"),
    (GREEN, "Rätt sak byggs första gången",
     "Specen är det kontrakt båda parter är överens om"),
    (GREEN, "Feedback sker i rätt fas",
     "Stakeholders granskar prototypen — inte färdig kod"),
    (GREEN, "Konsekvent slutresultat",
     "Designsystemet följs — inga enskilda tolkningar i koden"),
]

for i, (fc, title, sub) in enumerate(outcomes):
    y = ry + i * (rh + rg)
    add_rect(slide, rx, y, rw, rh, fill=GREEN_BG, border=GREEN_BD, border_width=Pt(1.5))
    add_text(slide, title, rx + Inches(0.2), y + Inches(0.1),  rw - Inches(0.3), Inches(0.32),
             size=Pt(13), bold=True, color=GREEN)
    add_text(slide, sub,   rx + Inches(0.2), y + Inches(0.44), rw - Inches(0.3), Inches(0.36),
             size=Pt(10), color=TEXT)

stmt_y = ry + 4 * (rh + rg) + Inches(0.08)
add_rect(slide, rx, stmt_y, rw, Inches(0.62), fill=BLUE_BG, border=BLUE_BD, border_width=Pt(1.5))
add_text(slide, "Design är inte en kostnad — det är det billigaste sättet att testa en idé.",
         rx + Inches(0.2), stmt_y + Inches(0.1), rw - Inches(0.4), Inches(0.45),
         size=Pt(12), bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────── SAVE
out = r"c:\Adrian\Code\txx–master-plan\_output\presentations\ui-ux-fokus.pptx"
prs.save(out)
print(f"Saved: {out}")
